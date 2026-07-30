"""
Extract plain text from stored files (PDF, Office, text, code, etc.).

Chat / file analysis can extract **any file size**:
- all PDF pages
- all Excel sheets and rows
- all PowerPoint slides
- full binary copy (no mid-file truncate)

Optional env caps (0 = unlimited): FILE_ANALYSIS_MAX_EXTRACT_CHARS
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
import tempfile
from pathlib import Path

from django.conf import settings

logger = logging.getLogger(__name__)

# Soft OCR trigger thresholds (not size caps)
OCR_TRIGGER_CHARS = 80
MIN_USEFUL_CHARS = 40

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".log",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".html",
    ".css",
    ".xml",
    ".yml",
    ".yaml",
    ".ini",
    ".cfg",
    ".env",
    ".docx",
    ".doc",
    ".rtf",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".pptx",
    ".ppt",
}

OFFICE_MIME_TYPES = {
    "application/pdf",
    "application/json",
    "application/xml",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-excel",
    "application/vnd.ms-powerpoint",
    "application/msword",
}

BINARY_DOC_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xlsm", ".xls", ".pptx", ".ppt"}


def _max_extract_chars() -> int:
    """0 or missing = keep full extracted text (any size)."""
    try:
        return int(getattr(settings, "FILE_ANALYSIS_MAX_EXTRACT_CHARS", 0) or 0)
    except (TypeError, ValueError):
        return 0


def extension_of(filename: str) -> str:
    return Path(filename or "").suffix.lower()


def is_analyzable(filename: str, mime_type: str = "") -> bool:
    ext = extension_of(filename)
    if ext in SUPPORTED_EXTENSIONS:
        return True
    mime = (mime_type or "").lower()
    return mime.startswith("text/") or mime in OFFICE_MIME_TYPES


def supported_types_label() -> str:
    return "PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), text, CSV, JSON, Markdown, and common code files"


def extract_text(file_obj, filename: str = "", mime_type: str = "") -> str:
    """
    Read from an open Django File / file-like object and return text.
    Prefer extract_from_storage_file() for Django FileField values.
    """
    ext = extension_of(filename)
    mime = (mime_type or "").lower()
    is_binary_doc = ext in BINARY_DOC_EXTENSIONS or any(
        token in mime
        for token in (
            "pdf",
            "wordprocessingml",
            "spreadsheetml",
            "presentationml",
            "ms-excel",
            "ms-powerpoint",
            "msword",
        )
    )

    if is_binary_doc:
        path = _materialize_temp(file_obj, suffix=ext or ".bin")
        if not path:
            return ""
        try:
            return extract_text_from_path(path, filename=filename, mime_type=mime_type)
        finally:
            try:
                path.unlink(missing_ok=True)
            except OSError:
                pass

    raw = _read_all_bytes(file_obj)
    if not raw:
        return ""
    return _extract_from_bytes(raw, filename=filename, mime_type=mime_type)


def extract_from_storage_file(storage_file, filename: str = "", mime_type: str = "") -> str:
    """
    Best entry point for Django FileField: uses local path when available
    so multi‑MB PDFs are not truncated in memory.
    """
    name = filename or getattr(storage_file, "name", "") or ""
    # Local filesystem storage exposes .path
    try:
        local_path = Path(storage_file.path)
        if local_path.exists():
            return extract_text_from_path(local_path, filename=name, mime_type=mime_type)
    except (AttributeError, NotImplementedError, ValueError, OSError):
        pass

    with storage_file.open("rb") as handle:
        return extract_text(handle, filename=name, mime_type=mime_type)


def extract_text_from_path(path: str | Path, filename: str = "", mime_type: str = "") -> str:
    path = Path(path)
    name = filename or path.name
    ext = extension_of(name)
    mime = (mime_type or "").lower()

    if ext == ".pdf" or mime == "application/pdf":
        return _from_pdf_path(path)
    if ext == ".docx" or "wordprocessingml" in mime:
        return _from_docx_path(path)
    if ext == ".doc":
        return (
            "Legacy .doc Word files are not supported. "
            "Please re-save the document as .docx and try again."
        )
    if ext in {".xlsx", ".xlsm"} or "spreadsheetml" in mime:
        return _from_xlsx_path(path)
    if ext == ".xls" or mime == "application/vnd.ms-excel":
        return _from_xls_path(path)
    if ext == ".pptx" or "presentationml" in mime:
        return _from_pptx_path(path)
    if ext == ".ppt":
        return (
            "Legacy .ppt PowerPoint files are not supported. "
            "Please re-save the deck as .pptx and try again."
        )

    # Full file — any size (text / CSV / JSON / code)
    raw = path.read_bytes()
    return _extract_from_bytes(raw, filename=name, mime_type=mime_type)


def _extract_from_bytes(raw: bytes, filename: str = "", mime_type: str = "") -> str:
    ext = extension_of(filename)
    mime = (mime_type or "").lower()

    if ext == ".csv":
        return _from_csv(raw, delimiter=",")
    if ext == ".tsv":
        return _from_csv(raw, delimiter="\t")
    if ext == ".json" or mime == "application/json":
        return _from_json(raw)
    if ext == ".rtf":
        return _from_rtf(raw)
    return _from_plain(raw)


def _materialize_temp(file_obj, suffix: str = ".bin") -> Path | None:
    """Copy the full file to a temp path (any size — no mid-file truncate)."""
    try:
        if hasattr(file_obj, "seek"):
            file_obj.seek(0)
        total = 0
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            target = Path(tmp.name)
            while True:
                chunk = file_obj.read(1024 * 1024)
                if not chunk:
                    break
                total += len(chunk)
                tmp.write(chunk)
        if hasattr(file_obj, "seek"):
            file_obj.seek(0)
        if total == 0:
            target.unlink(missing_ok=True)
            return None
        return target
    except Exception:
        logger.exception("Failed to materialize temp file for extraction")
        return None


def _read_all_bytes(file_obj) -> bytes:
    """Read the entire stream into memory (used for text-like formats)."""
    try:
        if hasattr(file_obj, "seek"):
            file_obj.seek(0)
        data = file_obj.read()
        if hasattr(file_obj, "seek"):
            file_obj.seek(0)
        return data or b""
    except Exception:
        return b""


def _from_plain(raw: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            text = ""
    else:
        text = raw.decode("utf-8", errors="ignore")
    return _clip(text)


def _from_pdf_path(path: Path) -> str:
    """
    1) Native text via PyMuPDF (best) or pypdf
    2) OCR on page images when the PDF is scanned / image-heavy
    """
    text = _pdf_text_pymupdf(path)
    if len(text.strip()) < OCR_TRIGGER_CHARS:
        ocr_text = _pdf_ocr_pymupdf(path)
        if len(ocr_text.strip()) > len(text.strip()):
            text = ocr_text
    if len(text.strip()) < OCR_TRIGGER_CHARS:
        ocr_text = _pdf_ocr_rapid(path)
        if len(ocr_text.strip()) > len(text.strip()):
            text = ocr_text
    if len(text.strip()) < MIN_USEFUL_CHARS:
        # Last resort: older pypdf path
        fallback = _pdf_text_pypdf(path)
        if len(fallback.strip()) > len(text.strip()):
            text = fallback
    return _clip(text)


def _pdf_text_pymupdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""

    try:
        doc = fitz.open(path)
    except Exception:
        return ""

    parts: list[str] = []
    try:
        # All pages — any PDF size
        for index in range(len(doc)):
            page = doc[index]
            # "blocks" preserves reading order better than raw text for reports.
            blocks = page.get_text("blocks") or []
            block_text = []
            for block in blocks:
                if len(block) >= 5 and isinstance(block[4], str):
                    block_text.append(block[4].strip())
            page_text = "\n".join(t for t in block_text if t) or (page.get_text("text") or "")
            parts.append(page_text)
    finally:
        doc.close()
    return "\n\n".join(parts).strip()


def _pdf_text_pypdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ""

    try:
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _pdf_ocr_pymupdf(path: Path) -> str:
    """OCR via PyMuPDF + Tesseract tessdata when installed on the machine."""
    try:
        import fitz
    except ImportError:
        return ""

    try:
        doc = fitz.open(path)
    except Exception:
        return ""

    parts: list[str] = []
    try:
        for index in range(len(doc)):
            page = doc[index]
            try:
                tp = page.get_textpage_ocr(language="eng", dpi=200, full=True)
                parts.append(page.get_text(textpage=tp) or "")
            except Exception as exc:
                logger.info("PyMuPDF OCR unavailable on page %s: %s", index, exc)
                break
    finally:
        doc.close()
    return "\n".join(parts).strip()


def _pdf_ocr_rapid(path: Path) -> str:
    """
    Pip-only OCR fallback (RapidOCR) — works for scanned/image PDFs
    without a system Tesseract install.
    """
    try:
        import fitz
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        return ""

    try:
        engine = RapidOCR()
        doc = fitz.open(path)
    except Exception:
        return ""

    parts: list[str] = []
    try:
        for index in range(len(doc)):
            page = doc[index]
            # ~150 DPI keeps OCR usable without huge memory for large prospectuses
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            png_bytes = pix.tobytes("png")
            try:
                result, _ = engine(png_bytes)
            except Exception as exc:
                logger.info("RapidOCR failed on page %s: %s", index, exc)
                continue
            if result:
                lines = [row[1] for row in result if len(row) > 1 and row[1]]
                if lines:
                    parts.append("\n".join(lines))
    finally:
        doc.close()
    return "\n\n".join(parts).strip()


def _from_docx_path(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""

    try:
        document = Document(str(path))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                parts.append(text)
        # Tables often hold requirements / data in Word reports.
        for table in document.tables:
            for row in table.rows:
                cells = [" ".join((cell.text or "").split()) for cell in row.cells]
                cells = [cell for cell in cells if cell]
                if cells:
                    parts.append(" | ".join(cells))
        return _clip("\n".join(parts))
    except Exception:
        logger.exception("DOCX extraction failed for %s", path)
        return ""


def _from_xlsx_path(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""

    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        logger.exception("XLSX open failed for %s", path)
        return ""

    parts: list[str] = []
    try:
        # All sheets and rows — any workbook size
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            parts.append(f"Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                values = [_stringify_cell(value) for value in row]
                values = [value for value in values if value]
                if not values:
                    continue
                parts.append(" | ".join(values))
    finally:
        workbook.close()
    return _clip("\n".join(parts))


def _from_xls_path(path: Path) -> str:
    """
    Legacy .xls support via xlrd when installed; otherwise ask user to convert.
    """
    try:
        import xlrd
    except ImportError:
        return (
            "Legacy .xls Excel files need the optional xlrd package, "
            "or re-save the workbook as .xlsx and try again."
        )

    try:
        book = xlrd.open_workbook(str(path))
        parts: list[str] = []
        for sheet_index in range(book.nsheets):
            sheet = book.sheet_by_index(sheet_index)
            parts.append(f"Sheet: {sheet.name}")
            for row_index in range(sheet.nrows):
                values = [_stringify_cell(sheet.cell_value(row_index, col)) for col in range(sheet.ncols)]
                values = [value for value in values if value]
                if values:
                    parts.append(" | ".join(values))
        return _clip("\n".join(parts))
    except Exception:
        logger.exception("XLS extraction failed for %s", path)
        return ""


def _from_pptx_path(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""

    try:
        deck = Presentation(str(path))
        parts: list[str] = []
        # All slides — any deck size
        for index, slide in enumerate(deck.slides):
            slide_lines: list[str] = [f"Slide {index + 1}"]
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    # Tables on slides
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [" ".join((cell.text or "").split()) for cell in row.cells]
                            cells = [cell for cell in cells if cell]
                            if cells:
                                slide_lines.append(" | ".join(cells))
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = " ".join((paragraph.text or "").split()).strip()
                    if text:
                        slide_lines.append(text)
            if len(slide_lines) > 1:
                parts.append("\n".join(slide_lines))
        return _clip("\n\n".join(parts))
    except Exception:
        logger.exception("PPTX extraction failed for %s", path)
        return ""


def _stringify_cell(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return f"{value:.6g}"
    text = str(value).strip()
    return text


def _from_docx(raw: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""

    try:
        document = Document(io.BytesIO(raw))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return _clip("\n".join(parts))
    except Exception:
        return ""


def _from_csv(raw: bytes, delimiter: str = ",") -> str:
    text = _from_plain(raw)
    try:
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        rows = []
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row))
        return _clip("\n".join(rows))
    except Exception:
        return text


def _from_json(raw: bytes) -> str:
    text = _from_plain(raw)
    try:
        data = json.loads(text)
        return _clip(json.dumps(data, indent=2, ensure_ascii=False))
    except Exception:
        return text


def _from_rtf(raw: bytes) -> str:
    text = _from_plain(raw)
    cleaned = re.sub(r"\\[a-zA-Z]+(-?\d+)?[ ]?", " ", text)
    cleaned = cleaned.replace("{", " ").replace("}", " ")
    return _clip(re.sub(r"\s+", " ", cleaned).strip())


def _clip(text: str) -> str:
    """Normalize whitespace; optionally clip only if settings set a max."""
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    limit = _max_extract_chars()
    if limit > 0 and len(text) > limit:
        return text[:limit] + "\n… [truncated for analysis]"
    return text
