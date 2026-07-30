"""
Tests for the local trainable file-analysis model and assistant integration.
"""

import io
import tempfile
from pathlib import Path
from unittest.mock import patch

from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import TestCase, override_settings

from accounts.models import Organization, User
from assistant.file_analysis.extractor import extract_text, extract_text_from_path
from assistant.file_analysis.intent import detect_intent
from assistant.file_analysis.model import SummarizerModel
from assistant.file_analysis.seed_corpus import iter_seed_texts
from assistant.services import AssistantService
from storage.models import FileNode


class SummarizerModelTests(TestCase):
    def test_train_and_summarize(self):
        docs = [text for _, text in iter_seed_texts()]
        model = SummarizerModel().train(docs, labels=["seed"] * len(docs))
        self.assertTrue(model.is_trained)
        self.assertGreater(model.vocabulary_size, 20)

        long_text = (
            "ABSTRACT This project builds a real-time fraud detection system for credit card transactions. "
            "Most credit card fraud detection systems rely on machine learning models trained on large "
            "transaction datasets. These models are good at spotting transactions that look fraudulent in "
            "general, but they share a blind spot because they judge each transaction on its own. "
            "A five hundred dollar purchase looks ordinary by itself, but it should raise an alarm if it "
            "happens far from where the cardholder lives on an unrecognized device. "
            "This project combines a machine learning layer with a behavioral context layer that scores "
            "every transaction against the individual cardholder history. "
            "The system is built with a Spring Boot backend, a Python scoring service, and a React dashboard. "
            "Evaluation shows the combined approach catches user-specific anomalies that a single model misses."
        )
        summary = model.summarize(long_text, max_sentences=3)
        self.assertTrue(summary)
        self.assertNotEqual(summary, "No readable text was found in this file.")
        self.assertIn("fraud", summary.lower())

    def test_keywords_and_question(self):
        docs = [text for _, text in iter_seed_texts()]
        model = SummarizerModel().train(docs)
        text = docs[1]
        keys = model.keywords(text, top_k=5)
        self.assertTrue(keys)
        answer = model.answer_question(text, "What are the remaining risks?")
        self.assertIn("risk", answer.lower())


class ExtractorTests(TestCase):
    def test_plain_and_json(self):
        from io import BytesIO

        plain = extract_text(BytesIO(b"Hello NexusStorage report."), filename="notes.txt")
        self.assertIn("Hello", plain)

        payload = b'{"title": "Demo", "count": 2}'
        dumped = extract_text(BytesIO(payload), filename="data.json")
        self.assertIn("Demo", dumped)

    def test_pdf_with_text_layer(self):
        try:
            import fitz
        except ImportError:
            self.skipTest("pymupdf not installed")

        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "NexusStorage prospectus revenue grew twelve percent.")
        pdf_bytes = doc.tobytes()
        doc.close()

        text = extract_text(io.BytesIO(pdf_bytes), filename="sample.pdf", mime_type="application/pdf")
        self.assertIn("prospectus", text.lower())

    def test_xlsx_and_pptx_extraction(self):
        try:
            from openpyxl import Workbook
            from pptx import Presentation
        except ImportError:
            self.skipTest("openpyxl/python-pptx not installed")

        from io import BytesIO
        from pathlib import Path
        import tempfile

        with tempfile.TemporaryDirectory() as tmp:
            xlsx_path = Path(tmp) / "budget.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.title = "Q1"
            ws.append(["Item", "Amount"])
            ws.append(["Cloud storage", 1200])
            ws.append(["Security audit", 800])
            wb.save(xlsx_path)
            xlsx_text = extract_text_from_path(xlsx_path, filename="budget.xlsx")
            self.assertIn("Cloud storage", xlsx_text)
            self.assertIn("1200", xlsx_text)

            pptx_path = Path(tmp) / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Fraud Detection Overview"
            slide.placeholders[1].text = "Behavioral scoring catches unusual spending patterns."
            prs.save(pptx_path)
            pptx_text = extract_text_from_path(pptx_path, filename="deck.pptx")
            self.assertIn("Fraud Detection", pptx_text)
            self.assertIn("Behavioral scoring", pptx_text)


class IntentTests(TestCase):
    def test_detect_summarize_pdf(self):
        intent = detect_intent("Please summarize quarterly-report.pdf")
        self.assertTrue(intent.wants_analysis)
        self.assertEqual(intent.filename_hint.lower(), "quarterly-report.pdf")
        self.assertEqual(intent.mode, "summary")

    def test_detect_know_about_pdf(self):
        intent = detect_intent("i want to know about propspectus2025.pdf")
        self.assertTrue(intent.wants_analysis)
        self.assertEqual(intent.filename_hint.lower(), "propspectus2025.pdf")
        self.assertEqual(intent.mode, "summary")

    def test_ignore_storage_question(self):
        intent = detect_intent("How much storage am I using?")
        self.assertFalse(intent.wants_analysis)


@override_settings(FILE_ANALYSIS_ENABLED=True, AI_PROVIDER="disabled")
class AssistantFileAnalysisTests(TestCase):
    def setUp(self):
        self.org = Organization.objects.create(name="Acme", slug="acme-file-ai")
        self.user = User.objects.create_user(
            email="analyst@acme.test",
            password="Passw0rd!",
            display_name="Analyst",
            organization=self.org,
        )
        docs = [text for _, text in iter_seed_texts()]
        self.model = SummarizerModel().train(docs, labels=["seed"] * len(docs))
        self._tmpdir = tempfile.TemporaryDirectory()
        self.model_path = Path(self._tmpdir.name) / "summarizer.json"
        self.model.save(self.model_path)

        body = (
            "Quarterly financial overview for Acme. Revenue grew twelve percent. "
            "Operating expenses remained stable. The board approved a new cloud storage budget. "
            "Risks include delayed vendor onboarding and incomplete audit evidence. "
            "Next quarter focuses on customer retention and security hardening."
        )
        upload = SimpleUploadedFile(
            "quarterly-report.txt",
            body.encode("utf-8"),
            content_type="text/plain",
        )
        self.node = FileNode.objects.create(
            organization=self.org,
            owner=self.user,
            name="quarterly-report.txt",
            node_type=FileNode.NodeType.FILE,
            content=upload,
            size_bytes=len(body),
            mime_type="text/plain",
        )

    def tearDown(self):
        self._tmpdir.cleanup()

    def test_summarize_named_file(self):
        with patch(
            "assistant.file_analysis.analyzer.SummarizerModel.load",
            return_value=SummarizerModel.load(self.model_path),
        ):
            service = AssistantService(self.user)
            answer, model = service.answer("summarize quarterly-report.txt", history=[])
        self.assertEqual(model, "nexus-file-analyzer")
        self.assertIn("quarterly-report.txt", answer)
        self.assertTrue(len(answer) > 40)

    def test_know_about_uses_file_contents(self):
        with patch(
            "assistant.file_analysis.analyzer.SummarizerModel.load",
            return_value=SummarizerModel.load(self.model_path),
        ):
            service = AssistantService(self.user)
            answer, model = service.answer("i want to know about quarterly-report.txt", history=[])
        self.assertEqual(model, "nexus-file-analyzer")
        self.assertIn("quarterly-report.txt", answer)
        self.assertNotIn("I can report storage usage", answer)

    def test_missing_file_message(self):
        with patch(
            "assistant.file_analysis.analyzer.SummarizerModel.load",
            return_value=SummarizerModel.load(self.model_path),
        ):
            service = AssistantService(self.user)
            answer, model = service.answer("summarize missing-file.pdf", history=[])
        self.assertEqual(model, "nexus-file-analyzer")
        self.assertIn("can't process", answer.lower())

    def test_unsupported_task_message(self):
        with patch(
            "assistant.file_analysis.analyzer.SummarizerModel.load",
            return_value=SummarizerModel.load(self.model_path),
        ):
            service = AssistantService(self.user)
            answer, model = service.answer("translate quarterly-report.txt to french", history=[])
        self.assertEqual(model, "nexus-file-analyzer")
        self.assertEqual(answer, "Sorry, I can't process this task.")
